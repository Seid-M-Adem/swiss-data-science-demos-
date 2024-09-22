import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.preprocessing import StandardScaler
from sklearn.decomposition import PCA
from sklearn.cluster import KMeans
from pptx import Presentation
from pptx.util import Inches
import os

# Paths
data_path = '/workspaces/swiss-data-science-demos-/customer_segmentation/data/customer_data.csv'
pptx_output_path = '/workspaces/swiss-data-science-demos-/customer_segmentation/Customer_Segmentation_Presentation.pptx'
graph_output_dir = '/workspaces/swiss-data-science-demos-/customer_segmentation/output/graphs'

# Create output directory for graphs if it doesn't exist
if not os.path.exists(graph_output_dir):
    os.makedirs(graph_output_dir)

# Load dataset
data = pd.read_csv(data_path)

# Select features for clustering
X = data[['age', 'average_spend', 'loyalty_score', 'engagement_score']]

# Standardize features
scaler = StandardScaler()
X_scaled = scaler.fit_transform(X)

# Apply K-means clustering
kmeans = KMeans(n_clusters=4, random_state=42)
data['cluster'] = kmeans.fit_predict(X_scaled)

# Reduce dimensions using PCA for visualization
pca = PCA(n_components=2)
X_pca = pca.fit_transform(X_scaled)

# Create a scatter plot of the clusters in PCA space
plt.figure(figsize=(10, 7))
sns.scatterplot(x=X_pca[:, 0], y=X_pca[:, 1], hue=data['cluster'], palette='Set1', s=100)
plt.title('Customer Segments (PCA Reduced)')
plt.xlabel('PCA 1')
plt.ylabel('PCA 2')
pca_graph_path = os.path.join(graph_output_dir, 'pca_clusters.png')
plt.savefig(pca_graph_path)
plt.close()

# Create PowerPoint presentation object
prs = Presentation()

# Function to add a slide with title and image to PowerPoint
def add_slide_with_image(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Add title slide
slide_layout = prs.slide_layouts[0]  # Title layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Customer Segmentation and Personalization"
subtitle.text = "An analysis using K-means clustering and PCA."

# Add PCA cluster visualization slide
add_slide_with_image(prs, 'Customer Segments (PCA Reduced)', pca_graph_path)

# Generate individual feature graphs
def create_and_save_graph(df, x_col, y_col, graph_filename, graph_title):
    plt.figure(figsize=(6, 4))
    sns.scatterplot(data=df, x=x_col, y=y_col, hue='cluster', palette='Set1', s=100)
    plt.title(graph_title)
    graph_file_path = os.path.join(graph_output_dir, graph_filename)
    plt.savefig(graph_file_path)
    plt.close()
    return graph_file_path

# Generate graphs for feature comparisons
graph1_path = create_and_save_graph(data, 'age', 'average_spend', 'age_vs_spend.png', 'Age vs Average Spend')
graph2_path = create_and_save_graph(data, 'loyalty_score', 'engagement_score', 'loyalty_vs_engagement.png', 'Loyalty Score vs Engagement Score')
graph3_path = create_and_save_graph(data, 'age', 'loyalty_score', 'age_vs_loyalty.png', 'Age vs Loyalty Score')
graph4_path = create_and_save_graph(data, 'age', 'engagement_score', 'age_vs_engagement.png', 'Age vs Engagement Score')

# Add graph slides to the presentation
add_slide_with_image(prs, 'Age vs Average Spend', graph1_path)
add_slide_with_image(prs, 'Loyalty Score vs Engagement Score', graph2_path)
add_slide_with_image(prs, 'Age vs Loyalty Score', graph3_path)
add_slide_with_image(prs, 'Age vs Engagement Score', graph4_path)

# Add conclusion slide
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
text_box = slide.shapes.placeholders[1]
text_box.text = (
    "We applied K-means clustering to segment customers into 4 distinct groups based on their age, "
    "spending behavior, loyalty, and engagement. PCA was used to reduce the dimensionality for visualization purposes. "
    "This analysis provides valuable insights into customer segmentation for targeted marketing."
)

# Save the presentation
prs.save(pptx_output_path)
print(f"PowerPoint presentation saved at {pptx_output_path}")
